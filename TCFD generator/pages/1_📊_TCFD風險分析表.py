#!/usr/bin/env python3
"""
TCFD 氣候風險分析表 - 互動式視覺化頁面
"""

import streamlit as st
import pandas as pd
import io
import os
from datetime import datetime
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 設定 output 資料夾路徑
OUTPUT_DIR = Path(__file__).parent.parent / "output"
OUTPUT_DIR.mkdir(exist_ok=True)


def create_tcfd_pptx():
    """生成大樓空調廠商 TCFD 風險分析 PowerPoint 簡報 (藍灰配色)"""
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # 16:9 寬螢幕
    prs.slide_height = Inches(7.5)
    
    # 顏色定義 - 藍灰配色
    BLUE_MAIN = RGBColor(74, 144, 164)    # #4a90a4
    GRAY_MAIN = RGBColor(122, 122, 122)   # #7a7a7a
    WHITE = RGBColor(255, 255, 255)
    LIGHT_GRAY = RGBColor(249, 249, 249)  # #f9f9f9
    
    # ========== 封面頁 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = BLUE_MAIN
    bg_shape.line.fill.background()
    
    # 裝飾區塊
    accent_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10), 0, Inches(3.33), Inches(7.5))
    accent_shape.fill.solid()
    accent_shape.fill.fore_color.rgb = GRAY_MAIN
    accent_shape.line.fill.background()
    
    # 標題
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "TCFD 氣候風險分析報告"
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # 副標題
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1.5))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "大樓空調系統廠商"
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(200, 230, 240)
    
    p2 = tf.add_paragraph()
    p2.text = "Task Force on Climate-related Financial Disclosures"
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(180, 210, 220)
    
    # 日期
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(9), Inches(0.5))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = datetime.now().strftime("%Y年%m月%d日")
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(180, 210, 220)
    
    # ========== TCFD 風險分析表 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 標題列
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(1.0))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = BLUE_MAIN
    title_bar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🌡️ TCFD 氣候風險分析表 - 大樓空調廠商"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # 資料內容
    tcfd_data = [
        ("極端高溫頻率增加\n夏季溫度持續上升，熱浪天數增加，\n導致冷卻需求大幅提升", 
         "設備負荷過重\n空調系統長時間高負荷運轉，\n設備壽命縮短，維修成本增加", 
         "開發高效能產品\n投資研發更高 EER 值的空調系統，\n提升極端氣候適應能力"),
        ("碳稅及環保法規\n政府實施碳稅制度，\n對高耗能設備課徵額外稅費", 
         "營運成本上升\n產品競爭力下降，\n客戶轉向選擇節能認證產品", 
         "取得綠色認證\n申請 ENERGY STAR、節能標章等認證，\n提升市場競爭力"),
        ("能源價格波動\n電力成本不穩定，再生能源需求增加，\n影響營運策略", 
         "客戶需求轉變\n大樓業主要求智能化節能方案，\n傳統產品需求下降", 
         "發展智慧空調系統\n整合 IoT 技術，\n提供 AI 控制及遠端監控功能")
    ]
    
    # 建立表格
    rows = len(tcfd_data) + 1
    cols = 3
    table = slide.shapes.add_table(rows, cols, Inches(0.3), Inches(1.2), Inches(12.73), Inches(5.8)).table
    
    # 設定欄寬
    table.columns[0].width = Inches(4.24)
    table.columns[1].width = Inches(4.24)
    table.columns[2].width = Inches(4.25)
    
    # 表頭
    headers = ["Description", "Impact", "Actions"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE_MAIN
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.size = Pt(18)
        para.font.color.rgb = WHITE
        para.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # 資料列
    for row_idx, (desc, impact, action) in enumerate(tcfd_data, 1):
        for col_idx, text in enumerate([desc, impact, action]):
            cell = table.cell(row_idx, col_idx)
            cell.text = text
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(11)
            para.alignment = PP_ALIGN.LEFT
            cell.vertical_anchor = MSO_ANCHOR.TOP
            
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
    
    # ========== 行動方案頁 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(1.0))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = BLUE_MAIN
    title_bar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "✅ 因應行動方案"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # 行動方案表格
    actions_data = [
        ("開發高效能產品", "投資 R&D 提升 EER 值", "2024-2025", "高"),
        ("取得綠色認證", "申請 ENERGY STAR 認證", "2024 Q2", "高"),
        ("發展智慧空調", "整合 IoT + AI 控制系統", "2024-2026", "中"),
        ("供應鏈減碳", "選用低碳原料供應商", "2025", "中"),
        ("員工培訓", "氣候風險意識教育", "持續進行", "低")
    ]
    
    table = slide.shapes.add_table(6, 4, Inches(0.5), Inches(1.3), Inches(12.33), Inches(5.5)).table
    
    table.columns[0].width = Inches(3.5)
    table.columns[1].width = Inches(4.5)
    table.columns[2].width = Inches(2)
    table.columns[3].width = Inches(2.33)
    
    headers = ["行動方案", "具體措施", "時程", "優先度"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE_MAIN
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.size = Pt(14)
        para.font.color.rgb = WHITE
        para.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    for row_idx, row_data in enumerate(actions_data, 1):
        for col_idx, text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = text
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(12)
            para.alignment = PP_ALIGN.CENTER if col_idx > 1 else PP_ALIGN.LEFT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
    
    # ========== 總結頁 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = BLUE_MAIN
    bg_shape.line.fill.background()
    
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(11), 0, Inches(2.33), Inches(7.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = GRAY_MAIN
    accent.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(10), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "📊 重點摘要"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    summary = [
        "🌡️ 極端高溫風險：開發高 EER 值空調系統",
        "💰 碳稅法規風險：取得 ENERGY STAR 等綠色認證",
        "⚡ 能源轉型風險：發展 IoT + AI 智慧空調",
        "🎯 策略目標：2025年前完成產品線升級",
        "📈 預期效益：提升市場競爭力 30%"
    ]
    
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(3), Inches(10), Inches(4))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(summary):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(22)
        p.font.color.rgb = WHITE
        p.space_after = Pt(16)
    
    # 備註
    note_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(10), Inches(0.5))
    tf = note_box.text_frame
    p = tf.paragraphs[0]
    p.text = "備註：此報告依據 TCFD 框架設計，建議定期檢視更新"
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(180, 210, 220)
    
    # 輸出
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output

st.set_page_config(
    page_title="TCFD 風險分析表",
    page_icon="📊",
    layout="wide"
)

# ============ 自定義 CSS ============
st.markdown("""
<style>
    @import url('https://fonts.googleapis.com/css2?family=Noto+Sans+TC:wght@300;400;500;700&display=swap');
    
    .tcfd-table {
        width: 100%;
        border-collapse: collapse;
        margin: 20px 0;
        font-family: 'Noto Sans TC', Arial, sans-serif;
        font-size: 14px;
    }

    .header-main {
        background: linear-gradient(135deg, #2d5a27 0%, #4a7c59 100%);
        color: white;
        text-align: center;
        padding: 16px;
        font-weight: bold;
        font-size: 18px;
        border-radius: 8px 8px 0 0;
    }

    .header-sub {
        background: linear-gradient(135deg, #f5f5f5 0%, #e8e8e8 100%);
        color: #333;
        text-align: center;
        padding: 12px;
        font-weight: bold;
        border: 1px solid #ddd;
    }

    .risk-category {
        background: #2d5a27;
        color: white;
        padding: 12px;
        font-weight: bold;
        text-align: center;
        vertical-align: middle;
    }

    .tcfd-table td {
        border: 1px solid #ddd;
        padding: 12px;
        vertical-align: top;
        line-height: 1.6;
    }

    .impact-high { 
        color: #d32f2f; 
        font-weight: bold; 
        background-color: #ffebee;
        padding: 2px 6px;
        border-radius: 4px;
    }
    
    .impact-medium { 
        color: #f57c00; 
        font-weight: bold;
        background-color: #fff3e0;
        padding: 2px 6px;
        border-radius: 4px;
    }
    
    .benefit { 
        color: #2e7d32; 
        font-weight: bold;
        background-color: #e8f5e9;
        padding: 2px 6px;
        border-radius: 4px;
    }

    .tech-highlight {
        background-color: #e3f2fd;
        padding: 3px 8px;
        border-radius: 4px;
        font-weight: 500;
        color: #1565c0;
    }
    
    .section-header {
        background: linear-gradient(135deg, #1a472a 0%, #2d5a27 100%);
        color: white;
        padding: 1rem 1.5rem;
        border-radius: 8px;
        margin: 1.5rem 0 1rem 0;
    }
</style>
""", unsafe_allow_html=True)

# ============ 標題 ============
st.markdown("""
<div class="section-header">
    <h2 style="margin:0; color:white;">📊 TCFD 氣候風險分析與節能減碳創新方案</h2>
</div>
""", unsafe_allow_html=True)

# ============ 風險分析表 ============
st.markdown("### 🌡️ 溫度上升對企業營運影響分析")

# ===== 大樓空調廠商 TCFD 風險數據 =====
hvac_risk_data = [
    {
        "description": "極端高溫頻率增加",
        "description_detail": "夏季溫度持續上升，熱浪天數增加，導致冷卻需求大幅提升",
        "impact": "設備負荷過重",
        "impact_detail": "空調系統長時間高負荷運轉，設備壽命縮短，維修成本增加",
        "actions": "開發高效能產品",
        "actions_detail": "投資研發更高 EER 值的空調系統，提升極端氣候適應能力"
    },
    {
        "description": "碳稅及環保法規",
        "description_detail": "政府實施碳稅制度，對高耗能設備課徵額外稅費",
        "impact": "營運成本上升",
        "impact_detail": "產品競爭力下降，客戶轉向選擇節能認證產品",
        "actions": "取得綠色認證",
        "actions_detail": "申請 ENERGY STAR、節能標章等認證，提升市場競爭力"
    },
    {
        "description": "能源價格波動",
        "description_detail": "電力成本不穩定，再生能源需求增加，影響營運策略",
        "impact": "客戶需求轉變",
        "impact_detail": "大樓業主要求智能化節能方案，傳統產品需求下降",
        "actions": "發展智慧空調系統",
        "actions_detail": "整合 IoT 技術，提供 AI 控制及遠端監控功能"
    }
]

# 舊版風險數據 (保留兼容)
risk_categories = ["大樓空調廠商 (新)", "設備營運風險", "員工健康風險", "能源供應風險"]
selected_category = st.selectbox("選擇風險類別", risk_categories)

# 風險數據
risk_data = {
    "設備營運風險": [
        {
            "描述": "極端高溫導致設備過熱當機",
            "影響": "🔴 停機損失每小時50-200萬<br>設備壽命減少15-25%",
            "措施": "部署 <span class='tech-highlight'>AI溫控系統</span><br>預防性維護，<span class='benefit'>降低30%故障率</span>"
        },
        {
            "描述": "冷卻系統能耗激增",
            "影響": "🟠 能源成本增加40-60%<br>碳排放量上升35%",
            "措施": "智能冷卻優化系統<br><span class='benefit'>節能25-40%</span>，ROI 2.5年"
        },
        {
            "描述": "戶外設施材料老化加速",
            "影響": "維護成本增加2-3倍<br>更換週期縮短50%",
            "措施": "採用耐候新材料<br>建立數位化巡檢系統"
        }
    ],
    "員工健康風險": [
        {
            "描述": "高溫作業環境健康風險",
            "影響": "🔴 中暑事故增加3倍<br>勞動生產力下降20%",
            "措施": "智能穿戴監測系統<br>動態調整作業時間"
        },
        {
            "描述": "室內空氣品質惡化",
            "影響": "員工請病假增加15%<br>工作效率降低12%",
            "措施": "<span class='tech-highlight'>AI空氣品質管理</span><br>即時調節通風系統"
        },
        {
            "描述": "通勤交通受極端天氣影響",
            "影響": "遲到缺勤率上升25%<br>營運連續性風險",
            "措施": "彈性工作制度<br>遠端辦公基礎設施"
        }
    ],
    "能源供應風險": [
        {
            "描述": "尖峰用電需求暴增",
            "影響": "🔴 電費支出增加50-80%<br>限電風險提高",
            "措施": "部署<span class='tech-highlight'>智能電網系統</span><br><span class='benefit'>削峰填谷30%</span>"
        },
        {
            "描述": "再生能源供應不穩定",
            "影響": "供電中斷風險增加<br>備用電源成本上升",
            "措施": "混合儲能系統<br>微電網建置，自給率達60%"
        },
        {
            "描述": "傳統能源價格波動",
            "影響": "🟠 能源成本波動±30%<br>預算規劃困難",
            "措施": "長期綠電採購合約<br>能源避險金融工具"
        }
    ]
}

# 顯示表格
def display_risk_table(category_data, category_name):
    st.markdown(f"#### {category_name}")
    
    table_html = """
    <table class="tcfd-table">
        <tr>
            <th class="header-sub" style="width:30%">風險描述</th>
            <th class="header-sub" style="width:35%">影響評估</th>
            <th class="header-sub" style="width:35%">適應措施</th>
        </tr>
    """
    
    for item in category_data:
        table_html += f"""
        <tr>
            <td>{item['描述']}</td>
            <td>{item['影響']}</td>
            <td>{item['措施']}</td>
        </tr>
        """
    
    table_html += "</table>"
    st.markdown(table_html, unsafe_allow_html=True)

if selected_category == "大樓空調廠商 (新)":
    # 顯示新的藍灰配色 HVAC TCFD 表格
    st.markdown("#### 🏢 大樓空調廠商 TCFD 氣候風險表")
    
    hvac_table_html = """
    <table style="width: 100%; border-collapse: collapse; font-family: Arial, sans-serif; margin: 20px 0;">
        <thead>
            <tr>
                <th style="background: linear-gradient(135deg, #4a90a4 50%, #7a7a7a 50%); color: white; padding: 15px; text-align: center; font-weight: bold; border: 1px solid #ddd; font-size: 16px;">
                    Description
                </th>
                <th style="background: linear-gradient(135deg, #4a90a4 50%, #7a7a7a 50%); color: white; padding: 15px; text-align: center; font-weight: bold; border: 1px solid #ddd; font-size: 16px;">
                    Impact
                </th>
                <th style="background: linear-gradient(135deg, #4a90a4 50%, #7a7a7a 50%); color: white; padding: 15px; text-align: center; font-weight: bold; border: 1px solid #ddd; font-size: 16px;">
                    Actions
                </th>
            </tr>
        </thead>
        <tbody>
    """
    
    for i, item in enumerate(hvac_risk_data):
        bg_color = "#f9f9f9" if i % 2 == 1 else "white"
        hvac_table_html += f"""
            <tr style="background-color: {bg_color};">
                <td style="padding: 12px; border: 1px solid #ddd; vertical-align: top;">
                    <strong>{item['description']}</strong><br>
                    {item['description_detail']}
                </td>
                <td style="padding: 12px; border: 1px solid #ddd; vertical-align: top;">
                    <strong>{item['impact']}</strong><br>
                    {item['impact_detail']}
                </td>
                <td style="padding: 12px; border: 1px solid #ddd; vertical-align: top;">
                    <strong>{item['actions']}</strong><br>
                    {item['actions_detail']}
                </td>
            </tr>
        """
    
    hvac_table_html += """
        </tbody>
    </table>
    <div style="margin-top: 10px; font-size: 12px; color: #666;">
        <strong>備註：</strong>此表格依據 TCFD 框架設計，協助大樓空調廠商識別氣候相關風險並制定對應策略。建議定期檢視更新內容，確保與最新氣候趨勢及法規要求同步。
    </div>
    """
    
    st.markdown(hvac_table_html, unsafe_allow_html=True)
else:
    for cat_name, cat_data in risk_data.items():
        if selected_category == cat_name:
            display_risk_table(cat_data, cat_name)
            break

st.markdown("---")

# ============ 創新節能減碳技術方案 ============
st.markdown("""
<div class="section-header">
    <h3 style="margin:0; color:white;">💡 創新節能減碳技術方案與效益</h3>
</div>
""", unsafe_allow_html=True)

# 解決方案數據
solution_df = pd.DataFrame({
    "技術方案": ["🤖 AI能耗監控系統", "🏗️ 被動式建築設計", "🏢 智能樓宇管理"],
    "技術特點": [
        "機器學習預測用電模式，即時優化設備運行參數",
        "自然通風、遮陽、保溫，減少機械空調依賴",
        "IoT感測整合控制，人員密度動態調節"
    ],
    "節能效益": ["15-25%", "30-40%", "25-35%"],
    "減碳效果": ["20%", "40%", "35%"],
    "投資回收期": ["1.8年", "3.5年", "2.2年"],
    "10年淨效益": ["+300萬", "+800萬", "+450萬"]
})

# 使用 Streamlit 原生表格顯示
st.dataframe(
    solution_df,
    use_container_width=True,
    hide_index=True,
    column_config={
        "技術方案": st.column_config.TextColumn("🔧 技術方案", width="medium"),
        "技術特點": st.column_config.TextColumn("📝 技術特點", width="large"),
        "節能效益": st.column_config.TextColumn("⚡ 節能效益", width="small"),
        "減碳效果": st.column_config.TextColumn("🌱 減碳效果", width="small"),
        "投資回收期": st.column_config.TextColumn("💰 ROI", width="small"),
        "10年淨效益": st.column_config.TextColumn("📈 10年淨效益", width="small")
    }
)

# ============ 互動式計算器 ============
st.markdown("---")
st.markdown("### 🧮 節能效益計算器")

col1, col2, col3 = st.columns(3)

with col1:
    current_energy_cost = st.number_input(
        "目前年度能源成本 (萬元)",
        min_value=0,
        max_value=100000,
        value=1000,
        step=100
    )

with col2:
    selected_tech = st.selectbox(
        "選擇導入技術",
        ["AI能耗監控系統 (節能20%)", "被動式建築設計 (節能35%)", "智能樓宇管理 (節能30%)"]
    )

with col3:
    carbon_price = st.number_input(
        "碳價格 (元/噸)",
        min_value=0,
        max_value=5000,
        value=500,
        step=50
    )

# 計算效益
tech_efficiency = {"AI能耗監控系統 (節能20%)": 0.20, "被動式建築設計 (節能35%)": 0.35, "智能樓宇管理 (節能30%)": 0.30}
efficiency = tech_efficiency[selected_tech]

energy_saving = current_energy_cost * efficiency
carbon_saving = energy_saving * 0.5 * carbon_price / 100  # 假設每萬元電費約0.5噸碳排

col1, col2, col3 = st.columns(3)

with col1:
    st.metric(
        label="💰 年度節省能源成本",
        value=f"{energy_saving:.0f} 萬元",
        delta=f"-{efficiency*100:.0f}%"
    )

with col2:
    st.metric(
        label="🌱 碳權價值估算",
        value=f"{carbon_saving:.1f} 萬元",
        delta="碳中和效益"
    )

with col3:
    total_benefit = energy_saving + carbon_saving
    st.metric(
        label="📈 總效益",
        value=f"{total_benefit:.0f} 萬元/年",
        delta="年化收益"
    )

# ============ 輸出報告區域 ============
st.markdown("---")
st.markdown("""
<div class="section-header">
    <h3 style="margin:0; color:white;">📥 輸出報告與下載</h3>
</div>
""", unsafe_allow_html=True)

# 風險數據 DataFrame
export_df = pd.DataFrame({
    "風險類別": ["設備", "設備", "設備", "員工", "員工", "員工", "能源", "能源", "能源"],
    "風險描述": ["設備過熱", "冷卻能耗", "材料老化", "健康風險", "空氣品質", "通勤影響", "尖峰用電", "供應不穩", "價格波動"],
    "影響描述": [
        "停機損失每小時50-200萬，設備壽命減少15-25%",
        "能源成本增加40-60%，碳排放量上升35%",
        "維護成本增加2-3倍，更換週期縮短50%",
        "中暑事故增加3倍，勞動生產力下降20%",
        "員工請病假增加15%，工作效率降低12%",
        "遲到缺勤率上升25%，營運連續性風險",
        "電費支出增加50-80%，限電風險提高",
        "供電中斷風險增加，備用電源成本上升",
        "能源成本波動±30%，預算規劃困難"
    ],
    "適應措施": [
        "部署AI溫控系統，預防性維護，降低30%故障率",
        "智能冷卻優化系統，節能25-40%，ROI 2.5年",
        "採用耐候新材料，建立數位化巡檢系統",
        "智能穿戴監測系統，動態調整作業時間",
        "AI空氣品質管理，即時調節通風系統",
        "彈性工作制度，遠端辦公基礎設施",
        "部署智能電網系統，削峰填谷30%",
        "混合儲能系統，微電網建置，自給率達60%",
        "長期綠電採購合約，能源避險金融工具"
    ],
    "影響程度": [9, 7, 6, 8, 5, 4, 9, 6, 7],
    "潛在損失(百萬)": [150, 80, 50, 30, 15, 10, 200, 100, 60]
})

# 節能方案 DataFrame
solution_export_df = pd.DataFrame({
    "技術方案": ["AI能耗監控系統", "被動式建築設計", "智能樓宇管理"],
    "技術特點": [
        "機器學習預測用電模式，即時優化設備運行參數",
        "自然通風、遮陽、保溫，減少機械空調依賴",
        "IoT感測整合控制，人員密度動態調節"
    ],
    "節能效益": ["15-25%", "30-40%", "25-35%"],
    "減碳效果": ["20%", "40%", "35%"],
    "投資回收期(年)": [1.8, 3.5, 2.2],
    "10年淨效益(萬)": [300, 800, 450]
})

# ===== 一鍵生成所有報告 =====
st.markdown("#### 🚀 一鍵生成所有報告")

col1, col2 = st.columns([1, 2])

with col1:
    generate_all = st.button("⚡ 生成所有報告到 output 資料夾", use_container_width=True, type="primary")

with col2:
    if generate_all:
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        saved_files = []
        
        with st.spinner("正在生成報告..."):
            # 1. 儲存 PPTX
            pptx_data = create_tcfd_pptx()
            pptx_path = OUTPUT_DIR / f"TCFD_報告_{timestamp}.pptx"
            with open(pptx_path, "wb") as f:
                f.write(pptx_data.getvalue())
            saved_files.append(f"✅ {pptx_path.name}")
            
            # 2. 儲存 CSV - 風險數據
            csv_path = OUTPUT_DIR / f"TCFD_風險數據_{timestamp}.csv"
            export_df.to_csv(csv_path, index=False, encoding="utf-8-sig")
            saved_files.append(f"✅ {csv_path.name}")
            
            # 3. 儲存 CSV - 節能方案
            solution_csv_path = OUTPUT_DIR / f"TCFD_節能方案_{timestamp}.csv"
            solution_export_df.to_csv(solution_csv_path, index=False, encoding="utf-8-sig")
            saved_files.append(f"✅ {solution_csv_path.name}")
            
            # 4. 儲存 Excel (包含多個工作表)
            excel_path = OUTPUT_DIR / f"TCFD_完整報告_{timestamp}.xlsx"
            with pd.ExcelWriter(excel_path, engine='openpyxl') as writer:
                export_df.to_excel(writer, sheet_name='風險分析', index=False)
                solution_export_df.to_excel(writer, sheet_name='節能方案', index=False)
            saved_files.append(f"✅ {excel_path.name}")
            
            # 5. 儲存 HTML
            html_path = OUTPUT_DIR / f"TCFD_風險表_{timestamp}.html"
            try:
                with open("TCFD/TCFD氣候風險表.py", "r", encoding="utf-8") as f:
                    html_content = f.read()
                with open(html_path, "w", encoding="utf-8") as f:
                    f.write(html_content)
                saved_files.append(f"✅ {html_path.name}")
            except:
                pass
        
        st.success(f"📁 已儲存 {len(saved_files)} 個檔案到 output 資料夾！")
        for f in saved_files:
            st.write(f)

st.markdown("---")

# ===== 個別下載按鈕 =====
st.markdown("#### 📁 個別下載")

col1, col2, col3, col4 = st.columns(4)

with col1:
    # 生成 PPTX
    pptx_data = create_tcfd_pptx()
    st.download_button(
        label="📽️ PowerPoint",
        data=pptx_data,
        file_name="TCFD_氣候風險分析報告.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        use_container_width=True
    )

with col2:
    # 生成 Excel
    excel_buffer = io.BytesIO()
    with pd.ExcelWriter(excel_buffer, engine='openpyxl') as writer:
        export_df.to_excel(writer, sheet_name='風險分析', index=False)
        solution_export_df.to_excel(writer, sheet_name='節能方案', index=False)
    excel_buffer.seek(0)
    
    st.download_button(
        label="📗 Excel 報告",
        data=excel_buffer,
        file_name="TCFD_完整報告.xlsx",
        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        use_container_width=True
    )

with col3:
    # 讀取 HTML 版本
    try:
        with open("TCFD/TCFD氣候風險表.py", "r", encoding="utf-8") as f:
            html_content = f.read()
        st.download_button(
            label="📄 HTML 網頁",
            data=html_content,
            file_name="TCFD_氣候風險分析表.html",
            mime="text/html",
            use_container_width=True
        )
    except:
        st.button("📄 HTML (無檔案)", disabled=True, use_container_width=True)

with col4:
    st.download_button(
        label="📊 CSV 數據",
        data=export_df.to_csv(index=False, encoding="utf-8-sig"),
        file_name="TCFD_風險數據.csv",
        mime="text/csv",
        use_container_width=True
    )

# ===== 儲存到本地按鈕 =====
st.markdown("---")
st.markdown("#### 💾 儲存到 output 資料夾")

col1, col2, col3, col4 = st.columns(4)

with col1:
    if st.button("💾 存 PPTX", use_container_width=True):
        pptx_data = create_tcfd_pptx()
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        pptx_path = OUTPUT_DIR / f"TCFD_報告_{timestamp}.pptx"
        with open(pptx_path, "wb") as f:
            f.write(pptx_data.getvalue())
        st.success(f"✅ 已儲存: {pptx_path.name}")

with col2:
    if st.button("💾 存 Excel", use_container_width=True):
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        excel_path = OUTPUT_DIR / f"TCFD_完整報告_{timestamp}.xlsx"
        with pd.ExcelWriter(excel_path, engine='openpyxl') as writer:
            export_df.to_excel(writer, sheet_name='風險分析', index=False)
            solution_export_df.to_excel(writer, sheet_name='節能方案', index=False)
        st.success(f"✅ 已儲存: {excel_path.name}")

with col3:
    if st.button("💾 存 HTML", use_container_width=True):
        try:
            with open("TCFD/TCFD氣候風險表.py", "r", encoding="utf-8") as f:
                html_content = f.read()
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            html_path = OUTPUT_DIR / f"TCFD_風險表_{timestamp}.html"
            with open(html_path, "w", encoding="utf-8") as f:
                f.write(html_content)
            st.success(f"✅ 已儲存: {html_path.name}")
        except:
            st.error("❌ HTML 來源檔案不存在")

with col4:
    if st.button("💾 存 CSV", use_container_width=True):
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        csv_path = OUTPUT_DIR / f"TCFD_風險數據_{timestamp}.csv"
        export_df.to_csv(csv_path, index=False, encoding="utf-8-sig")
        st.success(f"✅ 已儲存: {csv_path.name}")

# ===== 顯示 output 資料夾內容 =====
st.markdown("---")
st.markdown("#### 📂 output 資料夾內容")

if OUTPUT_DIR.exists():
    files = list(OUTPUT_DIR.glob("*"))
    if files:
        file_info = []
        for f in sorted(files, key=lambda x: x.stat().st_mtime, reverse=True):
            size_kb = f.stat().st_size / 1024
            mtime = datetime.fromtimestamp(f.stat().st_mtime).strftime("%Y-%m-%d %H:%M:%S")
            file_info.append({
                "檔案名稱": f.name,
                "大小": f"{size_kb:.1f} KB",
                "修改時間": mtime
            })
        
        st.dataframe(
            pd.DataFrame(file_info),
            use_container_width=True,
            hide_index=True
        )
        
        # 清空資料夾按鈕
        if st.button("🗑️ 清空 output 資料夾", type="secondary"):
            for f in files:
                f.unlink()
            st.success("✅ 已清空 output 資料夾")
            st.rerun()
    else:
        st.info("📭 output 資料夾是空的")
else:
    st.warning("⚠️ output 資料夾不存在")

# ============ 側邊欄 ============
with st.sidebar:
    st.markdown("### 🔗 快速連結")
    st.page_link("app.py", label="🏠 首頁")
    st.page_link("pages/1_📊_TCFD風險分析表.py", label="📊 TCFD 風險分析表")
    st.page_link("pages/2_🤖_Claude_AI助手.py", label="🤖 Claude AI 助手")
    st.page_link("pages/3_📈_數據分析工具.py", label="📈 數據分析工具")
    
    st.divider()
    
    st.markdown("### 📌 說明")
    st.info("""
    此頁面整合 TCFD 氣候風險框架，
    包含三大風險類別的完整分析。
    
    使用計算器可估算導入
    節能技術後的效益。
    """)


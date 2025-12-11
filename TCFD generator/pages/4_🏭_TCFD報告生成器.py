#!/usr/bin/env python3
"""
TCFD 報告生成器 - 輸入產業自動生成 PPTX
"""

import streamlit as st
import anthropic
import json
import io
import re
from datetime import datetime
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 設定 output 資料夾
OUTPUT_DIR = Path(__file__).parent.parent / "output"
OUTPUT_DIR.mkdir(exist_ok=True)

st.set_page_config(
    page_title="TCFD 報告生成器",
    page_icon="🏭",
    layout="wide"
)

# ============ 自定義樣式 ============
st.markdown("""
<style>
    .main-header {
        background: linear-gradient(135deg, #4a90a4 0%, #7a7a7a 100%);
        color: white;
        padding: 1.5rem;
        border-radius: 10px;
        text-align: center;
        margin-bottom: 2rem;
    }
    .success-box {
        background: #d4edda;
        border: 1px solid #c3e6cb;
        padding: 1rem;
        border-radius: 8px;
        margin: 1rem 0;
    }
    .info-box {
        background: #e7f3ff;
        border: 1px solid #b6d4fe;
        padding: 1rem;
        border-radius: 8px;
        margin: 1rem 0;
    }
</style>
""", unsafe_allow_html=True)

st.markdown("""
<div class="main-header">
    <h1 style="margin:0; color:white;">🏭 TCFD 報告生成器</h1>
    <p style="margin:0.5rem 0 0 0; opacity:0.9;">輸入您的產業，AI 自動生成 TCFD 氣候風險報告 + PPTX</p>
</div>
""", unsafe_allow_html=True)


# ============ PPTX 生成函數 ============
def create_industry_tcfd_pptx(industry_name, tcfd_data):
    """根據產業和 AI 生成的數據建立 PPTX"""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # 顏色
    BLUE_MAIN = RGBColor(74, 144, 164)
    GRAY_MAIN = RGBColor(122, 122, 122)
    WHITE = RGBColor(255, 255, 255)
    LIGHT_GRAY = RGBColor(249, 249, 249)
    
    # ========== 封面頁 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = BLUE_MAIN
    bg.line.fill.background()
    
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10), 0, Inches(3.33), Inches(7.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = GRAY_MAIN
    accent.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "TCFD 氣候風險分析報告"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = industry_name
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(200, 230, 240)
    
    p2 = tf.add_paragraph()
    p2.text = "Task Force on Climate-related Financial Disclosures"
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(180, 210, 220)
    
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(9), Inches(0.5))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = datetime.now().strftime("%Y年%m月%d日")
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(180, 210, 220)
    
    # ========== 風險分析表 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(1.0))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = BLUE_MAIN
    title_bar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"🌡️ TCFD 氣候風險分析 - {industry_name}"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # 表格
    risks = tcfd_data.get("risks", [])
    if risks:
        rows = len(risks) + 1
        table = slide.shapes.add_table(rows, 3, Inches(0.3), Inches(1.2), Inches(12.73), Inches(5.8)).table
        
        table.columns[0].width = Inches(4.24)
        table.columns[1].width = Inches(4.24)
        table.columns[2].width = Inches(4.25)
        
        headers = ["Description 風險描述", "Impact 影響評估", "Actions 因應措施"]
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = BLUE_MAIN
            para = cell.text_frame.paragraphs[0]
            para.font.bold = True
            para.font.size = Pt(14)
            para.font.color.rgb = WHITE
            para.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        for row_idx, risk in enumerate(risks, 1):
            for col_idx, key in enumerate(["description", "impact", "actions"]):
                cell = table.cell(row_idx, col_idx)
                cell.text = risk.get(key, "")
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(11)
                para.alignment = PP_ALIGN.LEFT
                cell.vertical_anchor = MSO_ANCHOR.TOP
                
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = LIGHT_GRAY
    
    # ========== 行動方案頁 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(1.0))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = BLUE_MAIN
    title_bar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "✅ 因應行動方案"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    actions = tcfd_data.get("action_plans", [])
    if actions:
        rows = len(actions) + 1
        table = slide.shapes.add_table(rows, 4, Inches(0.5), Inches(1.3), Inches(12.33), Inches(5.5)).table
        
        table.columns[0].width = Inches(3.5)
        table.columns[1].width = Inches(4.5)
        table.columns[2].width = Inches(2)
        table.columns[3].width = Inches(2.33)
        
        headers = ["行動方案", "具體措施", "時程", "優先度"]
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = BLUE_MAIN
            para = cell.text_frame.paragraphs[0]
            para.font.bold = True
            para.font.size = Pt(14)
            para.font.color.rgb = WHITE
            para.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        for row_idx, action in enumerate(actions, 1):
            data = [action.get("name", ""), action.get("measure", ""), 
                    action.get("timeline", ""), action.get("priority", "")]
            for col_idx, text in enumerate(data):
                cell = table.cell(row_idx, col_idx)
                cell.text = text
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(12)
                para.alignment = PP_ALIGN.CENTER if col_idx > 1 else PP_ALIGN.LEFT
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = LIGHT_GRAY
    
    # ========== 總結頁 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = BLUE_MAIN
    bg.line.fill.background()
    
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(11), 0, Inches(2.33), Inches(7.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = GRAY_MAIN
    accent.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(10), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "📊 重點摘要"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    summary = tcfd_data.get("summary", [])
    if summary:
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(10), Inches(4))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, item in enumerate(summary):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(20)
            p.font.color.rgb = WHITE
            p.space_after = Pt(12)
    
    # 備註
    note_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(10), Inches(0.5))
    tf = note_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"備註：此報告依據 TCFD 框架為{industry_name}設計，建議定期檢視更新"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(180, 210, 220)
    
    # 輸出
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output


def parse_ai_response(response_text):
    """解析 AI 回應，提取 TCFD 數據"""
    # 嘗試提取 JSON
    json_match = re.search(r'```json\s*([\s\S]*?)\s*```', response_text)
    if json_match:
        try:
            return json.loads(json_match.group(1))
        except:
            pass
    
    # 如果沒有 JSON，嘗試結構化解析
    tcfd_data = {
        "risks": [],
        "action_plans": [],
        "summary": []
    }
    
    # 簡單解析（按段落）
    lines = response_text.split('\n')
    current_section = None
    current_risk = {}
    
    for line in lines:
        line = line.strip()
        if not line:
            continue
            
        # 檢測風險段落
        if '風險' in line and ('描述' in line or 'Description' in line):
            current_section = 'description'
        elif '影響' in line or 'Impact' in line:
            current_section = 'impact'
        elif '措施' in line or '行動' in line or 'Action' in line:
            current_section = 'actions'
        elif line.startswith(('1.', '2.', '3.', '•', '-', '●')):
            # 新的項目
            if current_risk and all(k in current_risk for k in ['description', 'impact', 'actions']):
                tcfd_data["risks"].append(current_risk)
                current_risk = {}
            
            text = re.sub(r'^[0-9.\-•●\s]+', '', line)
            if current_section:
                current_risk[current_section] = text
    
    # 添加最後一個風險
    if current_risk:
        tcfd_data["risks"].append(current_risk)
    
    return tcfd_data


# ============ 側邊欄 ============
with st.sidebar:
    st.markdown("### 🔗 快速連結")
    st.page_link("app.py", label="🏠 首頁")
    st.page_link("pages/1_📊_TCFD風險分析表.py", label="📊 TCFD 風險分析表")
    st.page_link("pages/2_🤖_Claude_AI助手.py", label="🤖 Claude AI 助手")
    st.page_link("pages/3_📈_數據分析工具.py", label="📈 數據分析工具")
    st.page_link("pages/4_🏭_TCFD報告生成器.py", label="🏭 TCFD 報告生成器")
    
    st.divider()
    
    st.header("⚙️ API 設定")
    api_key = st.text_input("Claude API Key", type="password")
    
    model = st.selectbox(
        "模型",
        ["claude-sonnet-4-20250514", "claude-opus-4-20250514", "claude-sonnet-3-5-20241022"]
    )

# ============ 主要內容 ============
st.markdown("### 📝 步驟 1：輸入您的產業")

col1, col2 = st.columns([3, 1])

with col1:
    industry_input = st.text_input(
        "產業名稱",
        placeholder="例如：鋁建材業、空調設備業、太陽能產業...",
        help="請輸入您想分析的產業類型"
    )

with col2:
    industry_presets = st.selectbox(
        "或選擇預設",
        ["自訂", "鋁建材業", "大樓空調業", "鋼鐵業", "電子製造業", "營建業", "紡織業"]
    )

if industry_presets != "自訂":
    industry_input = industry_presets

# 生成按鈕
st.markdown("### 🚀 步驟 2：生成報告")

if st.button("⚡ 生成 TCFD 報告", type="primary", use_container_width=True):
    if not api_key:
        st.error("❌ 請先在側邊欄輸入 Claude API Key!")
    elif not industry_input:
        st.error("❌ 請輸入產業名稱!")
    else:
        with st.spinner(f"🤖 AI 正在分析 {industry_input} 的氣候風險..."):
            try:
                client = anthropic.Anthropic(api_key=api_key)
                
                prompt = f"""請為「{industry_input}」產業生成一份 TCFD 氣候風險分析報告。

請嚴格按照以下 JSON 格式輸出：

```json
{{
    "industry": "{industry_input}",
    "risks": [
        {{
            "description": "風險1標題\\n詳細描述...",
            "impact": "影響1標題\\n詳細影響...",
            "actions": "措施1標題\\n詳細措施..."
        }},
        {{
            "description": "風險2標題\\n詳細描述...",
            "impact": "影響2標題\\n詳細影響...",
            "actions": "措施2標題\\n詳細措施..."
        }},
        {{
            "description": "風險3標題\\n詳細描述...",
            "impact": "影響3標題\\n詳細影響...",
            "actions": "措施3標題\\n詳細措施..."
        }}
    ],
    "action_plans": [
        {{"name": "方案名稱1", "measure": "具體措施", "timeline": "2024-2025", "priority": "高"}},
        {{"name": "方案名稱2", "measure": "具體措施", "timeline": "2024-2026", "priority": "中"}},
        {{"name": "方案名稱3", "measure": "具體措施", "timeline": "2025", "priority": "中"}},
        {{"name": "方案名稱4", "measure": "具體措施", "timeline": "持續進行", "priority": "低"}}
    ],
    "summary": [
        "重點摘要1：關於主要風險",
        "重點摘要2：關於影響評估",
        "重點摘要3：關於因應策略",
        "重點摘要4：關於預期效益",
        "重點摘要5：關於時程目標"
    ]
}}
```

請確保：
1. risks 包含 3 個主要氣候風險項目
2. 每個風險都要有 description（風險描述）、impact（影響評估）、actions（因應措施）
3. action_plans 包含 4-5 個具體行動方案
4. summary 包含 5 個重點摘要
5. 內容要針對「{industry_input}」產業的特性撰寫
6. 只輸出 JSON，不要其他說明文字"""

                response = client.messages.create(
                    model=model,
                    max_tokens=4096,
                    temperature=0.3,
                    messages=[{"role": "user", "content": prompt}]
                )
                
                ai_response = response.content[0].text
                
                # 儲存到 session state
                st.session_state['ai_response'] = ai_response
                st.session_state['industry'] = industry_input
                
                # 解析 JSON
                json_match = re.search(r'```json\s*([\s\S]*?)\s*```', ai_response)
                if json_match:
                    tcfd_data = json.loads(json_match.group(1))
                    st.session_state['tcfd_data'] = tcfd_data
                    st.success("✅ AI 分析完成！請查看下方結果並下載報告")
                else:
                    # 嘗試直接解析
                    try:
                        tcfd_data = json.loads(ai_response)
                        st.session_state['tcfd_data'] = tcfd_data
                        st.success("✅ AI 分析完成！請查看下方結果並下載報告")
                    except:
                        st.warning("⚠️ AI 回應格式不完整，請查看原始回應")
                        st.session_state['tcfd_data'] = None
                
            except Exception as e:
                st.error(f"❌ API 錯誤: {e}")

# ============ 顯示結果 ============
if 'tcfd_data' in st.session_state and st.session_state.get('tcfd_data'):
    st.markdown("---")
    st.markdown("### 📊 步驟 3：查看與下載報告")
    
    tcfd_data = st.session_state['tcfd_data']
    industry = st.session_state.get('industry', '未知產業')
    
    # 顯示風險表格
    st.markdown(f"#### 🌡️ {industry} - TCFD 氣候風險分析")
    
    risks = tcfd_data.get("risks", [])
    if risks:
        # 建立 HTML 表格
        table_html = """
        <table style="width:100%; border-collapse:collapse; margin:1rem 0;">
            <thead>
                <tr>
                    <th style="background:linear-gradient(135deg,#4a90a4 50%,#7a7a7a 50%); color:white; padding:12px; border:1px solid #ddd;">Description</th>
                    <th style="background:linear-gradient(135deg,#4a90a4 50%,#7a7a7a 50%); color:white; padding:12px; border:1px solid #ddd;">Impact</th>
                    <th style="background:linear-gradient(135deg,#4a90a4 50%,#7a7a7a 50%); color:white; padding:12px; border:1px solid #ddd;">Actions</th>
                </tr>
            </thead>
            <tbody>
        """
        
        for i, risk in enumerate(risks):
            bg = "#f9f9f9" if i % 2 == 1 else "white"
            desc = risk.get("description", "").replace("\n", "<br>")
            impact = risk.get("impact", "").replace("\n", "<br>")
            actions = risk.get("actions", "").replace("\n", "<br>")
            
            table_html += f"""
                <tr style="background:{bg};">
                    <td style="padding:12px; border:1px solid #ddd; vertical-align:top;">{desc}</td>
                    <td style="padding:12px; border:1px solid #ddd; vertical-align:top;">{impact}</td>
                    <td style="padding:12px; border:1px solid #ddd; vertical-align:top;">{actions}</td>
                </tr>
            """
        
        table_html += "</tbody></table>"
        st.markdown(table_html, unsafe_allow_html=True)
    
    # 下載按鈕
    st.markdown("#### 📥 下載報告")
    
    col1, col2, col3 = st.columns(3)
    
    with col1:
        # 生成 PPTX
        pptx_data = create_industry_tcfd_pptx(industry, tcfd_data)
        st.download_button(
            label="📽️ 下載 PowerPoint",
            data=pptx_data,
            file_name=f"TCFD_{industry}_報告.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True
        )
    
    with col2:
        # 儲存到 output
        if st.button("💾 儲存到 output 資料夾", use_container_width=True):
            pptx_data = create_industry_tcfd_pptx(industry, tcfd_data)
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            pptx_path = OUTPUT_DIR / f"TCFD_{industry}_{timestamp}.pptx"
            with open(pptx_path, "wb") as f:
                f.write(pptx_data.getvalue())
            st.success(f"✅ 已儲存: {pptx_path.name}")
    
    with col3:
        # 下載 JSON
        st.download_button(
            label="📄 下載 JSON 數據",
            data=json.dumps(tcfd_data, ensure_ascii=False, indent=2),
            file_name=f"TCFD_{industry}_數據.json",
            mime="application/json",
            use_container_width=True
        )

# 顯示原始 AI 回應
if 'ai_response' in st.session_state:
    with st.expander("🔍 查看 AI 原始回應"):
        st.code(st.session_state['ai_response'], language="json")

# ============ 頁腳 ============
st.markdown("---")
st.caption("💡 提示：輸入產業名稱後，AI 會自動生成符合 TCFD 框架的氣候風險分析報告")



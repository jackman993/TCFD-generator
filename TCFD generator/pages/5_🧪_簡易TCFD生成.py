#!/usr/bin/env python3
"""
TCFD 報告生成器 (CSV 格式)
流程：用戶輸入產業 → LLM 回傳 CSV → 引擎製作 PPTX → 存到 output
"""

import streamlit as st
import anthropic
from datetime import datetime
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Output 路徑
OUTPUT_DIR = Path(__file__).parent.parent / "output"
OUTPUT_DIR.mkdir(exist_ok=True)

st.set_page_config(page_title="TCFD生成器", page_icon="🏭", layout="wide")

st.title("🏭 TCFD 報告生成器")
st.caption("輸入產業 → LLM 回傳 CSV → 引擎製作 PPTX")

# ============ 輸入區 ============
col1, col2 = st.columns(2)

with col1:
    api_key = st.text_input("Claude API Key", type="password")

with col2:
    industry = st.text_input("輸入您的產業", placeholder="例如：鋁建材業")

# ============ 生成按鈕 ============
if st.button("🚀 生成 TCFD 報告", type="primary", use_container_width=True):
    
    if not api_key:
        st.error("❌ 請輸入 API Key")
        st.stop()
    
    if not industry:
        st.error("❌ 請輸入產業")
        st.stop()
    
    # ========== Step 1: 呼叫 LLM，要求回傳 CSV ==========
    st.info("📡 Step 1: 呼叫 LLM...")
    
    prompt = f'''針對「{industry}」的 TCFD 氣候風險分析。

請輸出 3 行，每行用 ||| 分隔三個欄位（Description、Impact、Actions）：

風險描述|||影響評估|||因應措施
風險描述|||影響評估|||因應措施
風險描述|||影響評估|||因應措施

只輸出這 3 行，不要其他文字：'''

    try:
        client = anthropic.Anthropic(api_key=api_key)
        
        response = client.messages.create(
            model="claude-sonnet-4-20250514",
            max_tokens=1024,
            messages=[{"role": "user", "content": prompt}]
        )
        
        llm_response = response.content[0].text.strip()
        st.success("✅ Step 1 完成：LLM 已回應")
        
        # 顯示原始回應
        with st.expander("🔍 LLM 原始回應"):
            st.code(llm_response)
        
    except Exception as e:
        st.error(f"❌ API 錯誤: {e}")
        st.stop()
    
    # ========== Step 2: 解析 ||| 格式 ==========
    st.info("🔍 Step 2: 解析資料...")
    
    try:
        risks = []
        lines = llm_response.strip().split('\n')
        
        for line in lines:
            line = line.strip()
            if '|||' in line:
                parts = line.split('|||')
                if len(parts) >= 3:
                    risks.append({
                        'Description': parts[0].strip(),
                        'Impact': parts[1].strip(),
                        'Actions': parts[2].strip()
                    })
        
        st.success(f"✅ Step 2 完成：解析到 {len(risks)} 個風險項目")
        
        # 顯示解析結果
        with st.expander("📋 解析結果", expanded=True):
            for i, risk in enumerate(risks, 1):
                st.markdown(f"**風險 {i}**")
                st.write(f"- Description: {risk.get('Description', '')}")
                st.write(f"- Impact: {risk.get('Impact', '')}")
                st.write(f"- Actions: {risk.get('Actions', '')}")
                st.markdown("---")
        
        if len(risks) == 0:
            st.error("❌ 沒有解析到任何風險項目")
            st.stop()
        
    except Exception as e:
        st.error(f"❌ 解析失敗: {e}")
        st.code(llm_response)
        st.stop()
    
    # ========== Step 3: 引擎製作 PPTX ==========
    st.info("📽️ Step 3: 製作 PPTX...")
    
    try:
        # 建立簡報
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        # 顏色
        BLUE = RGBColor(74, 144, 164)
        GRAY = RGBColor(122, 122, 122)
        WHITE = RGBColor(255, 255, 255)
        LIGHT_BG = RGBColor(249, 249, 249)
        
        # ===== 封面頁 =====
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BLUE
        bg.line.fill.background()
        
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10), 0, Inches(3.33), Inches(7.5))
        accent.fill.solid()
        accent.fill.fore_color.rgb = GRAY
        accent.line.fill.background()
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "TCFD 氣候風險分析報告"
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = WHITE
        
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = industry
        p.font.size = Pt(32)
        p.font.color.rgb = RGBColor(200, 230, 240)
        
        date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(9), Inches(0.5))
        tf = date_box.text_frame
        p = tf.paragraphs[0]
        p.text = datetime.now().strftime("%Y年%m月%d日")
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(180, 210, 220)
        
        # ===== 表格頁 =====
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        header_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(1))
        header_bar.fill.solid()
        header_bar.fill.fore_color.rgb = BLUE
        header_bar.line.fill.background()
        
        header_text = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.6))
        tf = header_text.text_frame
        p = tf.paragraphs[0]
        p.text = f"TCFD 氣候風險分析 - {industry}"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = WHITE
        
        # 建立表格
        rows = len(risks) + 1
        cols = 3
        table = slide.shapes.add_table(rows, cols, Inches(0.3), Inches(1.2), Inches(12.73), Inches(5.8)).table
        
        table.columns[0].width = Inches(4.24)
        table.columns[1].width = Inches(4.24)
        table.columns[2].width = Inches(4.25)
        
        # 表頭
        headers = ["Description", "Impact", "Actions"]
        for i, h in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = BLUE
            para = cell.text_frame.paragraphs[0]
            para.font.bold = True
            para.font.size = Pt(16)
            para.font.color.rgb = WHITE
            para.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # 填入資料
        for row_idx, risk in enumerate(risks, 1):
            cell_data = [
                risk.get("Description", ""),
                risk.get("Impact", ""),
                risk.get("Actions", "")
            ]
            
            for col_idx, text in enumerate(cell_data):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(text)
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(11)
                para.alignment = PP_ALIGN.LEFT
                cell.vertical_anchor = MSO_ANCHOR.TOP
                
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = LIGHT_BG
        
        # ===== 儲存到 output =====
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"TCFD_{industry}_{timestamp}.pptx"
        filepath = OUTPUT_DIR / filename
        
        prs.save(str(filepath))
        
        st.success(f"✅ Step 3 完成：PPTX 已儲存")
        
    except Exception as e:
        st.error(f"❌ PPTX 製作失敗: {e}")
        import traceback
        st.code(traceback.format_exc())
        st.stop()
    
    # ========== 完成 ==========
    st.markdown("---")
    st.balloons()
    st.success(f"🎉 報告生成完成！")
    st.info(f"📁 檔案: `output/{filename}`")
    
    # 下載按鈕
    with open(filepath, "rb") as f:
        st.download_button(
            "📥 下載 PPTX",
            data=f.read(),
            file_name=filename,
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True
        )

# ============ 顯示 output 資料夾 ============
st.markdown("---")
st.markdown("### 📂 Output 資料夾")

files = sorted(OUTPUT_DIR.glob("*.pptx"), key=lambda x: x.stat().st_mtime, reverse=True)

if files:
    for f in files[:5]:
        col1, col2 = st.columns([4, 1])
        with col1:
            st.write(f"📄 {f.name}")
        with col2:
            st.write(f"{f.stat().st_size / 1024:.1f} KB")
else:
    st.info("尚無檔案")

if st.button("🔄 重新整理"):
    st.rerun()

#!/usr/bin/env python3
"""
Claude AI 助手頁面
整合 Claude API 的智能對話系統
AI 回答後自動生成 PPTX
"""

import streamlit as st
import anthropic
from pathlib import Path
import base64
from datetime import datetime
import json
import re
from docx import Document
import PyPDF2
from PIL import Image
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 設定 output 資料夾
OUTPUT_DIR = Path(__file__).parent.parent / "output"
OUTPUT_DIR.mkdir(exist_ok=True)

# ============ 頁面設定 ============
st.set_page_config(
    page_title="Claude AI 助手",
    page_icon="🤖",
    layout="wide"
)

st.title("🤖 Claude AI 智能助手")
st.caption("支援 TCFD 報告生成 → AI 回答後自動產生 PPTX")


# ============ PPTX 生成函數 ============
def parse_tcfd_from_response(response_text):
    """從 AI 回應中解析 TCFD 表格內容"""
    tcfd_items = []
    
    # 嘗試多種解析方式
    
    # 方式1: 尋找 HTML 表格
    table_match = re.search(r'<table[^>]*>(.*?)</table>', response_text, re.DOTALL | re.IGNORECASE)
    if table_match:
        table_html = table_match.group(1)
        # 提取 <tr> 中的 <td> 內容
        rows = re.findall(r'<tr[^>]*>(.*?)</tr>', table_html, re.DOTALL | re.IGNORECASE)
        for row in rows[1:]:  # 跳過表頭
            cells = re.findall(r'<td[^>]*>(.*?)</td>', row, re.DOTALL | re.IGNORECASE)
            if len(cells) >= 3:
                # 清理 HTML 標籤
                desc = re.sub(r'<[^>]+>', '\n', cells[0]).strip()
                impact = re.sub(r'<[^>]+>', '\n', cells[1]).strip()
                actions = re.sub(r'<[^>]+>', '\n', cells[2]).strip()
                tcfd_items.append({
                    "description": desc,
                    "impact": impact,
                    "actions": actions
                })
    
    # 方式2: 尋找 Markdown 表格
    if not tcfd_items:
        md_rows = re.findall(r'\|([^|]+)\|([^|]+)\|([^|]+)\|', response_text)
        for row in md_rows:
            if '---' not in row[0] and 'Description' not in row[0] and '描述' not in row[0]:
                tcfd_items.append({
                    "description": row[0].strip(),
                    "impact": row[1].strip(),
                    "actions": row[2].strip()
                })
    
    # 方式3: 尋找編號列表
    if not tcfd_items:
        # 嘗試找 **標題** 格式
        sections = re.split(r'\n(?=\d+\.|\*\*|###)', response_text)
        current_item = {}
        
        for section in sections:
            section = section.strip()
            if not section:
                continue
            
            lower_section = section.lower()
            
            if 'description' in lower_section or '風險描述' in lower_section or '描述' in lower_section:
                if current_item and 'description' in current_item:
                    tcfd_items.append(current_item)
                    current_item = {}
                # 提取內容
                content = re.sub(r'^[\d\.\*\#\s]+', '', section)
                content = re.sub(r'\*\*[^*]+\*\*', '', content, count=1).strip()
                current_item['description'] = content[:200] if content else section[:200]
                
            elif 'impact' in lower_section or '影響' in lower_section:
                content = re.sub(r'^[\d\.\*\#\s]+', '', section)
                content = re.sub(r'\*\*[^*]+\*\*', '', content, count=1).strip()
                current_item['impact'] = content[:200] if content else section[:200]
                
            elif 'action' in lower_section or '措施' in lower_section or '行動' in lower_section:
                content = re.sub(r'^[\d\.\*\#\s]+', '', section)
                content = re.sub(r'\*\*[^*]+\*\*', '', content, count=1).strip()
                current_item['actions'] = content[:200] if content else section[:200]
        
        if current_item and len(current_item) >= 2:
            tcfd_items.append(current_item)
    
    return tcfd_items


def create_tcfd_pptx_from_response(industry_name, tcfd_items, full_response):
    """根據 AI 回應建立 PPTX"""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # 顏色 - 藍灰配色
    BLUE_MAIN = RGBColor(74, 144, 164)
    GRAY_MAIN = RGBColor(122, 122, 122)
    WHITE = RGBColor(255, 255, 255)
    LIGHT_GRAY = RGBColor(249, 249, 249)
    
    # ========== 封面頁 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = BLUE_MAIN
    bg.line.fill.background()
    
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10), 0, Inches(3.33), Inches(7.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = GRAY_MAIN
    accent.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "TCFD 氣候風險分析報告"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = industry_name
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(200, 230, 240)
    
    p2 = tf.add_paragraph()
    p2.text = "Task Force on Climate-related Financial Disclosures"
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(180, 210, 220)
    
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(9), Inches(0.5))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = datetime.now().strftime("%Y年%m月%d日")
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(180, 210, 220)
    
    # ========== TCFD 表格頁 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(1.0))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = BLUE_MAIN
    title_bar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"🌡️ TCFD 氣候風險分析 - {industry_name}"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # 如果有解析到 TCFD 項目，建立表格
    if tcfd_items:
        rows = len(tcfd_items) + 1
        table = slide.shapes.add_table(rows, 3, Inches(0.3), Inches(1.2), Inches(12.73), Inches(5.8)).table
        
        table.columns[0].width = Inches(4.24)
        table.columns[1].width = Inches(4.24)
        table.columns[2].width = Inches(4.25)
        
        # 表頭
        headers = ["Description", "Impact", "Actions"]
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = BLUE_MAIN
            para = cell.text_frame.paragraphs[0]
            para.font.bold = True
            para.font.size = Pt(16)
            para.font.color.rgb = WHITE
            para.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # 資料列
        for row_idx, item in enumerate(tcfd_items, 1):
            for col_idx, key in enumerate(["description", "impact", "actions"]):
                cell = table.cell(row_idx, col_idx)
                text = item.get(key, "")
                # 限制長度避免超出
                cell.text = text[:300] if len(text) > 300 else text
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(10)
                para.alignment = PP_ALIGN.LEFT
                cell.vertical_anchor = MSO_ANCHOR.TOP
                
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = LIGHT_GRAY
    else:
        # 如果沒解析到，顯示完整回應
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.33), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        # 限制長度
        display_text = full_response[:2000] + "..." if len(full_response) > 2000 else full_response
        p.text = display_text
        p.font.size = Pt(11)
    
    # ========== 總結頁 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = BLUE_MAIN
    bg.line.fill.background()
    
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(11), 0, Inches(2.33), Inches(7.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = GRAY_MAIN
    accent.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(10), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "📊 報告摘要"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    summary_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(10), Inches(4))
    tf = summary_box.text_frame
    tf.word_wrap = True
    
    if tcfd_items:
        p = tf.paragraphs[0]
        p.text = f"• 產業：{industry_name}"
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE
        
        p2 = tf.add_paragraph()
        p2.text = f"• 識別風險項目：{len(tcfd_items)} 項"
        p2.font.size = Pt(20)
        p2.font.color.rgb = WHITE
        
        p3 = tf.add_paragraph()
        p3.text = "• 依據 TCFD 框架分析"
        p3.font.size = Pt(20)
        p3.font.color.rgb = WHITE
        
        p4 = tf.add_paragraph()
        p4.text = f"• 報告生成時間：{datetime.now().strftime('%Y-%m-%d %H:%M')}"
        p4.font.size = Pt(20)
        p4.font.color.rgb = WHITE
    else:
        p = tf.paragraphs[0]
        p.text = "AI 回應內容已匯入報告"
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE
    
    # 備註
    note_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(10), Inches(0.5))
    tf = note_box.text_frame
    p = tf.paragraphs[0]
    p.text = "此報告由 AI 自動生成，建議專業審閱後使用"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(180, 210, 220)
    
    # 輸出
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output


def extract_industry_from_messages(messages):
    """從對話中提取產業名稱"""
    for msg in reversed(messages):
        content = msg.get("content", "")
        if isinstance(content, list):
            content = " ".join([c.get("text", "") for c in content if c.get("type") == "text"])
        
        # 尋找產業關鍵字
        match = re.search(r'(我是|我們是|屬於|從事|經營)?[「「]?([^「」\s]{2,10}(?:業|產業|公司|廠商|製造|工業))[」」]?', content)
        if match:
            return match.group(2)
        
        # 常見產業
        industries = ["鋁建材", "空調", "鋼鐵", "電子", "紡織", "營建", "太陽能", "半導體", "汽車", "化工", "食品"]
        for ind in industries:
            if ind in content:
                return ind + "業"
    
    return "企業"


# ============ 側邊欄設定 ============
with st.sidebar:
    st.markdown("### 🔗 快速連結")
    st.page_link("app.py", label="🏠 首頁")
    st.page_link("pages/1_📊_TCFD風險分析表.py", label="📊 TCFD 風險分析表")
    st.page_link("pages/2_🤖_Claude_AI助手.py", label="🤖 Claude AI 助手")
    st.page_link("pages/3_📈_數據分析工具.py", label="📈 數據分析工具")
    
    st.divider()
    
    st.header("⚙️ API 設定")
    
    api_key = st.text_input(
        "Claude API Key",
        type="password",
        help="從 https://console.anthropic.com 取得"
    )
    
    model = st.selectbox(
        "模型選擇",
        ["claude-sonnet-4-20250514", "claude-opus-4-20250514", "claude-sonnet-3-5-20241022"]
    )
    
    st.subheader("🎛️ 參數")
    max_tokens = st.slider("Max Tokens", 1024, 8192, 4096)
    temperature = st.slider("Temperature", 0.0, 1.0, 0.5, 0.1)
    
    st.divider()
    
    # 自動生成 PPTX 開關
    st.subheader("📽️ PPTX 設定")
    auto_generate_pptx = st.checkbox("AI 回答後自動生成 PPTX", value=True)
    
    st.divider()
    
    col1, col2 = st.columns(2)
    with col1:
        if st.button("🗑️ 清除對話", use_container_width=True):
            st.session_state.messages = []
            st.session_state.total_cost = 0
            st.session_state.last_pptx = None
            st.rerun()
    
    with col2:
        if st.button("💾 下載對話", use_container_width=True):
            if 'messages' in st.session_state:
                chat_json = json.dumps(st.session_state.messages, indent=2, ensure_ascii=False)
                st.download_button(
                    "下載 JSON",
                    chat_json,
                    file_name=f"chat_{datetime.now().strftime('%Y%m%d_%H%M%S')}.json",
                    mime="application/json"
                )
    
    st.divider()
    st.subheader("💰 成本統計")
    if 'total_cost' not in st.session_state:
        st.session_state.total_cost = 0
    st.metric("本次總成本", f"${st.session_state.total_cost:.4f}")


# ============ 初始化 Session State ============
if 'messages' not in st.session_state:
    st.session_state.messages = []

if 'client' not in st.session_state:
    st.session_state.client = None

if 'last_pptx' not in st.session_state:
    st.session_state.last_pptx = None

if 'pending_template' not in st.session_state:
    st.session_state.pending_template = None


# ============ 工具函數 ============
def read_file_content(file):
    try:
        file_type = file.name.split('.')[-1].lower()
        if file_type == 'txt':
            return file.read().decode('utf-8')
        elif file_type == 'docx':
            doc = Document(io.BytesIO(file.read()))
            return '\n'.join([para.text for para in doc.paragraphs])
        elif file_type == 'pdf':
            pdf_reader = PyPDF2.PdfReader(io.BytesIO(file.read()))
            return ''.join([page.extract_text() for page in pdf_reader.pages])
    except:
        return None

def encode_image(image_file):
    try:
        image = Image.open(image_file)
        buffered = io.BytesIO()
        image.save(buffered, format=image.format or "PNG")
        mime_type = f"image/{image.format.lower()}" if image.format else "image/png"
        return {
            "type": "image",
            "source": {
                "type": "base64",
                "media_type": mime_type,
                "data": base64.b64encode(buffered.getvalue()).decode('utf-8')
            }
        }
    except:
        return None

def calculate_cost(input_tokens, output_tokens, model_name):
    if "sonnet" in model_name.lower():
        return input_tokens / 1_000_000 * 3 + output_tokens / 1_000_000 * 15
    elif "opus" in model_name.lower():
        return input_tokens / 1_000_000 * 15 + output_tokens / 1_000_000 * 75
    return 0


# ============ TCFD 快捷模板 ============
st.markdown("### ⚡ TCFD 報告生成")
st.info("💡 點擊下方按鈕，再輸入您的產業（如：我是鋁建材業），AI 回答後會自動生成 PPTX")

col1, col2, col3 = st.columns(3)

with col1:
    if st.button("📊 生成 TCFD 報告書", use_container_width=True, type="primary"):
        st.session_state.pending_template = """請為我的產業分析 3 個主要氣候風險。

輸出格式必須是純 CSV，不要輸出 HTML，不要輸出 Markdown。
格式如下：
Description,Impact,Actions
風險描述1,影響評估1,因應措施1
風險描述2,影響評估2,因應措施2
風險描述3,影響評估3,因應措施3

請針對我的產業特性撰寫具體、專業的內容。現在請直接輸出 CSV："""
        st.rerun()

with col2:
    if st.button("🌡️ 溫度風險分析", use_container_width=True):
        st.session_state.pending_template = """分析溫度上升對我的產業營運影響，輸出 3 項風險。
格式必須是純 CSV：
Description,Impact,Actions
風險1,影響1,措施1
風險2,影響2,措施2
風險3,影響3,措施3
不要輸出 HTML，直接輸出 CSV："""
        st.rerun()

with col3:
    if st.button("💡 節能方案建議", use_container_width=True):
        st.session_state.pending_template = """請為我的產業提供 3 個節能減碳方案。
格式必須是純 CSV：
Description,Impact,Actions
方案1描述,方案1效益,方案1措施
方案2描述,方案2效益,方案2措施
方案3描述,方案3效益,方案3措施
不要輸出 HTML，直接輸出 CSV："""
        st.rerun()

# 顯示待發送的模板
if st.session_state.pending_template:
    st.warning("📝 模板已準備好！請在下方輸入您的產業後送出")
    with st.expander("查看模板內容"):
        st.code(st.session_state.pending_template)


# ============ 顯示對話歷史 ============
st.markdown("---")
st.markdown("### 💬 對話")

for message in st.session_state.messages:
    with st.chat_message(message["role"]):
        if isinstance(message["content"], str):
            st.markdown(message["content"])
        elif isinstance(message["content"], list):
            for block in message["content"]:
                if block.get("type") == "text":
                    st.markdown(block.get("text", ""))


# ============ 顯示上次生成的 PPTX ============
if st.session_state.last_pptx:
    st.markdown("---")
    st.markdown("### 📥 下載報告")
    
    col1, col2, col3 = st.columns(3)
    
    with col1:
        st.download_button(
            "📽️ 下載 PowerPoint",
            data=st.session_state.last_pptx['data'],
            file_name=st.session_state.last_pptx['filename'],
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True
        )
    
    with col2:
        if st.button("💾 儲存到 output", use_container_width=True):
            pptx_path = OUTPUT_DIR / st.session_state.last_pptx['filename']
            with open(pptx_path, "wb") as f:
                f.write(st.session_state.last_pptx['data'].getvalue())
            st.success(f"✅ 已儲存: {pptx_path.name}")
    
    with col3:
        st.caption(f"🏭 產業: {st.session_state.last_pptx.get('industry', '未知')}")
        st.caption(f"📊 風險項目: {st.session_state.last_pptx.get('items_count', 0)} 項")


# ============ 用戶輸入 ============
user_input = st.chat_input("輸入您的產業（如：我是鋁建材業）...")

if user_input:
    if not api_key:
        st.error("❌ 請先在側邊欄輸入 Claude API Key!")
        st.stop()
    
    if st.session_state.client is None:
        st.session_state.client = anthropic.Anthropic(api_key=api_key)
    
    # 組合訊息：模板 + 用戶輸入
    if st.session_state.pending_template:
        full_message = f"{user_input}\n\n{st.session_state.pending_template}"
        st.session_state.pending_template = None
    else:
        full_message = user_input
    
    # 保存用戶訊息
    st.session_state.messages.append({
        "role": "user",
        "content": full_message
    })
    
    with st.chat_message("user"):
        st.markdown(user_input)
        if "TCFD" in full_message:
            st.caption("📋 已附加 TCFD 報告模板")
    
    # 調用 Claude API
    with st.chat_message("assistant"):
        with st.spinner("🤔 AI 分析中..."):
            try:
                system_prompt = """你是專業的 TCFD 氣候風險顧問。
請用繁體中文回答。
當被要求生成 TCFD 表格時，請務必使用 HTML <table> 格式輸出，包含完整的 <tr><td> 標籤。
每個風險項目要具體針對用戶的產業特性撰寫。"""

                response = st.session_state.client.messages.create(
                    model=model,
                    max_tokens=max_tokens,
                    temperature=temperature,
                    system=system_prompt,
                    messages=st.session_state.messages
                )
                
                assistant_message = response.content[0].text
                st.markdown(assistant_message)
                
                # 計算成本
                cost = calculate_cost(response.usage.input_tokens, response.usage.output_tokens, model)
                st.session_state.total_cost += cost
                st.caption(f"📊 Tokens: {response.usage.input_tokens} in / {response.usage.output_tokens} out | 💰 ${cost:.4f}")
                
                # 保存助手訊息
                st.session_state.messages.append({
                    "role": "assistant",
                    "content": assistant_message
                })
                
                # ====== 自動生成 PPTX ======
                if auto_generate_pptx:
                    with st.spinner("📽️ 正在生成 PPTX..."):
                        # 提取產業名稱
                        industry = extract_industry_from_messages(st.session_state.messages)
                        
                        # 解析 TCFD 內容
                        tcfd_items = parse_tcfd_from_response(assistant_message)
                        
                        # 生成 PPTX
                        pptx_data = create_tcfd_pptx_from_response(industry, tcfd_items, assistant_message)
                        
                        # 儲存到 session state
                        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                        filename = f"TCFD_{industry}_{timestamp}.pptx"
                        
                        st.session_state.last_pptx = {
                            'data': pptx_data,
                            'filename': filename,
                            'industry': industry,
                            'items_count': len(tcfd_items)
                        }
                        
                        # 自動儲存到 output
                        pptx_path = OUTPUT_DIR / filename
                        with open(pptx_path, "wb") as f:
                            f.write(pptx_data.getvalue())
                        
                        st.success(f"✅ PPTX 已自動生成並儲存到 output/{filename}")
                        st.info(f"📊 解析到 {len(tcfd_items)} 個風險項目")
                
            except Exception as e:
                st.error(f"❌ 錯誤: {e}")


# ============ 頁腳 ============
st.divider()
st.caption("💡 流程：點擊「生成 TCFD 報告書」→ 輸入產業 → AI 回答 → 自動生成 PPTX 到 output 資料夾")

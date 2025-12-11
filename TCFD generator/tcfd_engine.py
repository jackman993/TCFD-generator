"""
TCFD 表格引擎 - 從 CSV 生成 PPTX
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path

OUTPUT_DIR = Path(__file__).parent / "output"
OUTPUT_DIR.mkdir(exist_ok=True)

RISK_TYPES = ['政策與法規', '綠色產品與科技']


def create_tcfd_table(csv_lines, industry="企業", filename=None):
    """從 CSV 生成 TCFD PPTX"""
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 標題
    title = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(12), Inches(0.5))
    p = title.text_frame.paragraphs[0]
    p.text = f"TCFD 轉型風險分析 - {industry}"
    p.font.size = Pt(20)
    p.font.bold = True
    
    # 表格
    rows = 2 + len(csv_lines)
    tbl = slide.shapes.add_table(rows, 6, Inches(0.15), Inches(0.85), Inches(13), Inches(6.2)).table
    
    # 欄寬（風險描述右邊界=中線6.5"）
    tbl.columns[0].width = Inches(1.5)   # 類型
    tbl.columns[1].width = Inches(1.8)   # 氣候風險
    tbl.columns[2].width = Inches(1.2)   # 期間
    tbl.columns[3].width = Inches(2.0)   # 風險描述（右邊界=6.5"）
    tbl.columns[4].width = Inches(3.2)   # 潛在影響
    tbl.columns[5].width = Inches(3.3)   # 因應措施
    
    # 表頭列高（+60%）
    tbl.rows[0].height = Inches(0.56)
    tbl.rows[1].height = Inches(0.56)
    
    # Row 0: 分割表頭（合併 cell 去除內部格線）
    # 左半：合併 0-2 欄
    tbl.cell(0, 0).merge(tbl.cell(0, 2))
    left_cell = tbl.cell(0, 0)
    left_cell.fill.solid()
    left_cell.fill.fore_color.rgb = RGBColor(0x2F, 0x52, 0x33)
    left_cell.text = "氣候相關風險"
    left_cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    left_cell.text_frame.paragraphs[0].font.bold = True
    left_cell.text_frame.paragraphs[0].font.size = Pt(11)
    left_cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 右半：合併 3-5 欄
    tbl.cell(0, 3).merge(tbl.cell(0, 5))
    right_cell = tbl.cell(0, 3)
    right_cell.fill.solid()
    right_cell.fill.fore_color.rgb = RGBColor(0x80, 0x80, 0x80)
    right_cell.text = "財務影響"
    right_cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    right_cell.text_frame.paragraphs[0].font.bold = True
    right_cell.text_frame.paragraphs[0].font.size = Pt(11)
    right_cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Row 1: 欄位標題
    headers = ['類型', '氣候風險', '期間', '風險描述', '潛在影響', '因應措施']
    for c, h in enumerate(headers):
        cell = tbl.cell(1, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x80, 0x80, 0x80)
        cell.text = h
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(10)
    
    # 資料行
    for i, line in enumerate(csv_lines):
        r = i + 2
        parts = [p.strip() for p in line.split('|||')]
        
        # 類型
        cell = tbl.cell(r, 0)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x8B, 0x9D, 0x83)
        if i == 0:
            cell.text = "轉型風險"
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(11)
        
        # 氣候風險
        risk = RISK_TYPES[i] if i < len(RISK_TYPES) else f"風險{i+1}"
        tbl.cell(r, 1).text = risk
        tbl.cell(r, 1).text_frame.paragraphs[0].font.size = Pt(11)
        
        # 期間
        tbl.cell(r, 2).text = "中短期"
        tbl.cell(r, 2).text_frame.paragraphs[0].font.size = Pt(11)
        
        # 描述、影響、措施（3點用分號分隔，轉成換行）
        if len(parts) >= 1:
            _set_bullet_text(tbl.cell(r, 3), parts[0])
        if len(parts) >= 2:
            _set_bullet_text(tbl.cell(r, 4), parts[1])
        if len(parts) >= 3:
            _set_bullet_text(tbl.cell(r, 5), parts[2])
    
    # 儲存
    if filename is None:
        from datetime import datetime
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"TCFD_{industry}_{timestamp}.pptx"
    
    filepath = OUTPUT_DIR / filename
    prs.save(str(filepath))
    return filepath


def _set_bullet_text(cell, text):
    """設定多點文字（用分號分隔轉成多行）"""
    tf = cell.text_frame
    tf.clear()
    
    # 用分號分隔成多點
    points = [p.strip() for p in text.split(';') if p.strip()]
    
    for idx, point in enumerate(points):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = f"• {point}"
        p.font.size = Pt(11)
        p.alignment = PP_ALIGN.LEFT


if __name__ == "__main__":
    test = [
        "碳排放申報義務加重;碳費徵收時程提前;合規要求趨嚴,年增120萬合規成本;需增聘ESG人員;罰款風險增加,建立碳管理系統;導入自動監測;定期內部稽核",
        "低碳製程壓力;客戶要求認證;技術門檻提高,資本支出增800萬;研發成本上升;毛利率下降,分階段導入;申請補助;技術合作",
        "綠色需求提升;傳統產品衰退;市場結構改變,營收轉型壓力;客戶流失風險;價格競爭,開發低碳產品;建立認證;拓展新市場"
    ]
    print(f"已生成：{create_tcfd_table(test, '測試')}")

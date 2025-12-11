#!/usr/bin/env python3
"""
TCFD 氣候風險 PPTX 生成器
針對大樓空調廠商設計
"""

import io
import os
from datetime import datetime
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 設定 output 資料夾路徑
OUTPUT_DIR = Path(__file__).parent / "output"
OUTPUT_DIR.mkdir(exist_ok=True)

def create_hvac_tcfd_pptx():
    """生成大樓空調廠商 TCFD 風險分析 PowerPoint 簡報"""
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # 16:9 寬螢幕
    prs.slide_height = Inches(7.5)
    
    # 顏色定義 - 使用藍灰漸變配色
    BLUE_MAIN = RGBColor(74, 144, 164)    # #4a90a4
    GRAY_MAIN = RGBColor(122, 122, 122)   # #7a7a7a
    WHITE = RGBColor(255, 255, 255)
    LIGHT_GRAY = RGBColor(249, 249, 249)  # #f9f9f9
    
    # ========== 封面頁 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = BLUE_MAIN
    bg_shape.line.fill.background()
    
    # 裝飾區塊
    accent_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10), 0, Inches(3.33), Inches(7.5))
    accent_shape.fill.solid()
    accent_shape.fill.fore_color.rgb = GRAY_MAIN
    accent_shape.line.fill.background()
    
    # 標題
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "TCFD 氣候風險分析報告"
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # 副標題
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "大樓空調系統廠商"
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(200, 230, 240)
    
    p2 = tf.add_paragraph()
    p2.text = "Task Force on Climate-related Financial Disclosures"
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(180, 210, 220)
    
    # 日期
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(9), Inches(0.5))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = datetime.now().strftime("%Y年%m月%d日")
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(180, 210, 220)
    
    # ========== TCFD 風險分析表 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 標題列
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(1.0))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = BLUE_MAIN
    title_bar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🌡️ TCFD 氣候風險分析表 - 大樓空調廠商"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # 資料內容
    tcfd_data = [
        {
            "description": "極端高溫頻率增加\n夏季溫度持續上升，熱浪天數增加，\n導致冷卻需求大幅提升",
            "impact": "設備負荷過重\n空調系統長時間高負荷運轉，\n設備壽命縮短，維修成本增加",
            "actions": "開發高效能產品\n投資研發更高 EER 值的空調系統，\n提升極端氣候適應能力"
        },
        {
            "description": "碳稅及環保法規\n政府實施碳稅制度，\n對高耗能設備課徵額外稅費",
            "impact": "營運成本上升\n產品競爭力下降，\n客戶轉向選擇節能認證產品",
            "actions": "取得綠色認證\n申請 ENERGY STAR、節能標章等認證，\n提升市場競爭力"
        },
        {
            "description": "能源價格波動\n電力成本不穩定，再生能源需求增加，\n影響營運策略",
            "impact": "客戶需求轉變\n大樓業主要求智能化節能方案，\n傳統產品需求下降",
            "actions": "發展智慧空調系統\n整合 IoT 技術，\n提供 AI 控制及遠端監控功能"
        }
    ]
    
    # 建立表格
    rows = len(tcfd_data) + 1  # 資料 + 表頭
    cols = 3
    left = Inches(0.3)
    top = Inches(1.2)
    width = Inches(12.73)
    height = Inches(5.8)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # 設定欄寬
    table.columns[0].width = Inches(4.24)
    table.columns[1].width = Inches(4.24)
    table.columns[2].width = Inches(4.25)
    
    # 表頭
    headers = ["Description", "Impact", "Actions"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE_MAIN
        
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.size = Pt(18)
        para.font.color.rgb = WHITE
        para.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # 資料列
    for row_idx, data in enumerate(tcfd_data, 1):
        for col_idx, key in enumerate(["description", "impact", "actions"]):
            cell = table.cell(row_idx, col_idx)
            cell.text = data[key]
            
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(12)
            para.alignment = PP_ALIGN.LEFT
            cell.vertical_anchor = MSO_ANCHOR.TOP
            
            # 設定第一行粗體
            if cell.text_frame.paragraphs:
                first_line = cell.text.split('\n')[0]
                cell.text_frame.paragraphs[0].runs[0].font.bold = True if cell.text_frame.paragraphs[0].runs else False
            
            # 交替背景色
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
    
    # ========== 風險詳細分析頁 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 標題
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(1.0))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = BLUE_MAIN
    title_bar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🔍 風險影響評估"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # 三個風險卡片
    risks = [
        ("🌡️ 極端高溫風險", "設備負荷增加 40%\n維修成本上升 25%\n產品壽命縮短 3-5 年", BLUE_MAIN),
        ("💰 法規合規風險", "碳稅成本增加\n市場准入門檻提高\n認證費用支出", GRAY_MAIN),
        ("⚡ 能源轉型風險", "客戶需求轉變\n傳統產品淘汰\n技術升級壓力", BLUE_MAIN)
    ]
    
    card_width = Inches(4)
    card_height = Inches(4.5)
    start_left = Inches(0.5)
    card_top = Inches(1.5)
    gap = Inches(0.33)
    
    for i, (title, content, color) in enumerate(risks):
        left_pos = start_left + i * (card_width + gap)
        
        # 卡片背景
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, card_top, card_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(245, 248, 250)
        card.line.color.rgb = color
        card.line.width = Pt(2)
        
        # 卡片標題
        card_title = slide.shapes.add_textbox(left_pos + Inches(0.2), card_top + Inches(0.2), card_width - Inches(0.4), Inches(0.6))
        tf = card_title.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = color
        
        # 卡片內容
        card_content = slide.shapes.add_textbox(left_pos + Inches(0.2), card_top + Inches(1), card_width - Inches(0.4), Inches(3))
        tf = card_content.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = content
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(80, 80, 80)
    
    # ========== 行動方案頁 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(1.0))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = BLUE_MAIN
    title_bar.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "✅ 因應行動方案"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # 行動方案表格
    actions_data = [
        ("開發高效能產品", "投資 R&D 提升 EER 值", "2024-2025", "高"),
        ("取得綠色認證", "申請 ENERGY STAR 認證", "2024 Q2", "高"),
        ("發展智慧空調", "整合 IoT + AI 控制系統", "2024-2026", "中"),
        ("供應鏈減碳", "選用低碳原料供應商", "2025", "中"),
        ("員工培訓", "氣候風險意識教育", "持續進行", "低")
    ]
    
    table = slide.shapes.add_table(6, 4, Inches(0.5), Inches(1.3), Inches(12.33), Inches(5.5)).table
    
    table.columns[0].width = Inches(3.5)
    table.columns[1].width = Inches(4.5)
    table.columns[2].width = Inches(2)
    table.columns[3].width = Inches(2.33)
    
    # 表頭
    headers = ["行動方案", "具體措施", "時程", "優先度"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE_MAIN
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.size = Pt(14)
        para.font.color.rgb = WHITE
        para.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    for row_idx, row_data in enumerate(actions_data, 1):
        for col_idx, text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = text
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(12)
            para.alignment = PP_ALIGN.CENTER if col_idx > 1 else PP_ALIGN.LEFT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
    
    # ========== 總結頁 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = BLUE_MAIN
    bg_shape.line.fill.background()
    
    # 裝飾
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(11), 0, Inches(2.33), Inches(7.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = GRAY_MAIN
    accent.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(10), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "📊 重點摘要"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    summary = [
        "🌡️ 極端高溫風險：開發高 EER 值空調系統",
        "💰 碳稅法規風險：取得 ENERGY STAR 等綠色認證",
        "⚡ 能源轉型風險：發展 IoT + AI 智慧空調",
        "🎯 策略目標：2025年前完成產品線升級",
        "📈 預期效益：提升市場競爭力 30%"
    ]
    
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(3), Inches(10), Inches(4))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(summary):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(22)
        p.font.color.rgb = WHITE
        p.space_after = Pt(16)
    
    # 備註
    note_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(10), Inches(0.5))
    tf = note_box.text_frame
    p = tf.paragraphs[0]
    p.text = "備註：此報告依據 TCFD 框架設計，建議定期檢視更新"
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(180, 210, 220)
    
    return prs


def save_to_output():
    """儲存 PPTX 到 output 資料夾"""
    prs = create_hvac_tcfd_pptx()
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    
    # 儲存 PPTX
    pptx_path = OUTPUT_DIR / f"TCFD_空調廠商報告_{timestamp}.pptx"
    prs.save(str(pptx_path))
    print(f"✅ 已儲存: {pptx_path}")
    
    return pptx_path


def get_pptx_bytes():
    """取得 PPTX 的 bytes 格式 (供 Streamlit 下載)"""
    prs = create_hvac_tcfd_pptx()
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output


if __name__ == "__main__":
    print("🚀 生成 TCFD 空調廠商報告...")
    saved_path = save_to_output()
    print(f"📁 檔案位置: {saved_path}")
    print("✅ 完成!")


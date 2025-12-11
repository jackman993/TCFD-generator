"""
TCFD 表格引擎 05 - 資源效率（機會）
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path

OUTPUT_DIR = Path(__file__).parent.parent / "output"
OUTPUT_DIR.mkdir(exist_ok=True)

TABLE_TITLE = "TCFD 資源效率分析"
TYPE_NAME = "機會面"
RISK_TYPES = ['能源效率提升', '資源循環利用']


def create_table(csv_lines, industry="企業", filename=None):
    """從 CSV 生成 TCFD PPTX"""
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 標題
    title = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(12), Inches(0.5))
    p = title.text_frame.paragraphs[0]
    p.text = f"{TABLE_TITLE} - {industry}"
    p.font.size = Pt(20)
    p.font.bold = True
    
    # 表格
    rows = 2 + len(csv_lines)
    tbl = slide.shapes.add_table(rows, 6, Inches(0.15), Inches(0.85), Inches(13), Inches(6.2)).table
    
    # 欄寬
    tbl.columns[0].width = Inches(1.5)
    tbl.columns[1].width = Inches(1.8)
    tbl.columns[2].width = Inches(1.2)
    tbl.columns[3].width = Inches(2.0)
    tbl.columns[4].width = Inches(3.2)
    tbl.columns[5].width = Inches(3.3)
    
    # 表頭列高
    tbl.rows[0].height = Inches(0.56)
    tbl.rows[1].height = Inches(0.56)
    
    # Row 0: 分割表頭（機會面用藍色）
    tbl.cell(0, 0).merge(tbl.cell(0, 2))
    left_cell = tbl.cell(0, 0)
    left_cell.fill.solid()
    left_cell.fill.fore_color.rgb = RGBColor(0x1F, 0x4E, 0x79)  # 藍色
    left_cell.text = "氣候相關機會"
    left_cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    left_cell.text_frame.paragraphs[0].font.bold = True
    left_cell.text_frame.paragraphs[0].font.size = Pt(11)
    left_cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    tbl.cell(0, 3).merge(tbl.cell(0, 5))
    right_cell = tbl.cell(0, 3)
    right_cell.fill.solid()
    right_cell.fill.fore_color.rgb = RGBColor(0x80, 0x80, 0x80)
    right_cell.text = "財務影響"
    right_cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    right_cell.text_frame.paragraphs[0].font.bold = True
    right_cell.text_frame.paragraphs[0].font.size = Pt(11)
    right_cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Row 1: 欄位標題
    headers = ['類型', '機會項目', '期間', '機會描述', '潛在效益', '行動方案']
    for c, h in enumerate(headers):
        cell = tbl.cell(1, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x80, 0x80, 0x80)
        cell.text = h
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(10)
    
    # 資料行
    for i, line in enumerate(csv_lines):
        r = i + 2
        parts = [p.strip() for p in line.split('|||')]
        
        cell = tbl.cell(r, 0)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x5B, 0x9B, 0xD5)  # 淺藍色
        if i == 0:
            cell.text = TYPE_NAME
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(11)
        
        risk = RISK_TYPES[i] if i < len(RISK_TYPES) else f"機會{i+1}"
        tbl.cell(r, 1).text = risk
        tbl.cell(r, 1).text_frame.paragraphs[0].font.size = Pt(11)
        
        tbl.cell(r, 2).text = "中長期"
        tbl.cell(r, 2).text_frame.paragraphs[0].font.size = Pt(11)
        
        if len(parts) >= 1:
            _set_bullet_text(tbl.cell(r, 3), parts[0])
        if len(parts) >= 2:
            _set_bullet_text(tbl.cell(r, 4), parts[1])
        if len(parts) >= 3:
            _set_bullet_text(tbl.cell(r, 5), parts[2])
    
    # 儲存
    if filename is None:
        from datetime import datetime
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"TCFD_05_資源效率_{industry}_{timestamp}.pptx"
    
    filepath = OUTPUT_DIR / filename
    prs.save(str(filepath))
    return filepath


def _set_bullet_text(cell, text):
    tf = cell.text_frame
    tf.clear()
    points = [p.strip() for p in text.split(';') if p.strip()]
    for idx, point in enumerate(points):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(11)
        p.alignment = PP_ALIGN.LEFT

